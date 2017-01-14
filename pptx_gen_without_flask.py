import sys
import zipfile
import StringIO
from pptx import Presentation
from pptx.util import Cm
from PIL import Image
import os
from flask import Flask, request, redirect, url_for, flash, render_template, make_response
from werkzeug.utils import secure_filename

# filename structuur binnen zipfile:
# meerdere coupletten:
# projectie-111-muziek-couplet-1-1.png
# maar 1 couplet:
# projectie-425-muziek-2.png
def song_couplets2arr(filenamelist):
	song_couplets = {}
	for filename in filenamelist:
		if filename[-3:] == 'png':
			title_text_arr = filename.split('-')
			#print filename
			#print title_text_arr[3][0]
			if title_text_arr[3][0].isdigit() == False:  # meerdere coupletten
				if title_text_arr[1] in song_couplets:
					if title_text_arr[4] not in song_couplets[title_text_arr[1]]:
						song_couplets[title_text_arr[1]].append(title_text_arr[4])
				else:
					song_couplets[title_text_arr[1]] = [title_text_arr[4]]
			else:  # 1 couplet
				song_couplets[title_text_arr[1]] = ['1']
	song_couplets[title_text_arr[1]].sort(key=int)
	return song_couplets

def get_song_title_text(filename, song_couplets):
	title_text = 'Lied '
	title_text_arr = filename.split('-')
	if title_text_arr[1].isdigit() and int(title_text_arr[1]) <= 150:
		title_text = "Psalm "
	title_text += title_text_arr[1] + ': '

	#print len(song_couplets[title_text_arr[1]])
	for couplet in song_couplets[title_text_arr[1]]:
		if len(song_couplets[title_text_arr[1]])==1:
			title_text += ' [' + couplet + '] '
		elif couplet == title_text_arr[4]:
			title_text += ' [' + couplet + '] '
		else:
			title_text += ' ' + couplet + ' '
	return title_text

def create_pptx(pptx_template_file):
	prs = Presentation(pptx_template_file)  # if filename is give, load the presentation
	prs.core_properties.author = "pptx-generator v0.1"
	prs.core_properties.title = "pptx-generator generator powerpoint for Hervgemb"
	return prs

def create_title_slide(prs, titel_tekst, sub_titel_tekst):
	title_slide_layout = prs.slide_layouts[0]  # layout 0 = de startpagina
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = titel_tekst
	subtitle.text = sub_titel_tekst

def create_song_slide(prs, song_title, song_img_data):
	song_slide_layout = prs.slide_layouts[1]  # layout 1 = titel + object
	slide = prs.slides.add_slide(song_slide_layout)

	# set title Psalm / Lied
	title = slide.shapes.title
	title.text = song_title
	#title.left = Cm(3.30)
	#title.top = Cm(1.50)

	# set song image (SongText + MusicNotes)
	left = Cm(3.29) 
	top = Cm(1.94)
	#pic = slide.shapes.add_picture(song_img_data, left, top)
	img4= StringIO.StringIO(song_img_data.getvalue())
	pic = slide.shapes.add_picture(img4, left, top)
	#pic = slide.shapes.add_picture(song_img_data, left, top)
	pic.height=pic.height/3
	pic.width=pic.width/3


def create_scripture_slide(prs, scripture_title, scripture_text):
	scripture_slide_layout = prs.slide_layouts[1]  # layout 1 = titel + object
	slide = prs.slides.add_slide(scripture_slide_layout)

	# set scripture title
	title = slide.shapes.title
	title.text = scripture_title

	# set scripture text
	subtitle = slide.placeholders[1]
	subtitle.top = Cm(1.94)
	subtitle.left = Cm(1.00)
	subtitle.width = prs.slide_width - subtitle.left - Cm(0.50)
	subtitle.height = prs.slide_height - subtitle.top - Cm(0.50)
	subtitle.text += scripture_text + '\n'

def create_intermediate_slide(prs, tekst):
	interim_slide_layout = prs.slide_layouts[3]  # layout 3 = intermediate_slide
	slide = prs.slides.add_slide(interim_slide_layout)

	# set title
	title = slide.shapes.title
	title.text = tekst
	#title.left = Cm(3.30)
	#title.top = Cm(1.50)

def create_index_slide(prs, song_couplets, scripture_fragments, datum_tekst):
	index_slide_layout = prs.slide_layouts[1]  # layout 1 = titel + object
	slide = prs.slides.add_slide(index_slide_layout)

	# set title Psalm / Lied
	title = slide.shapes.title
	title.text = 'Liturgie ' + datum_tekst
	#title.left = Cm(3.30)
	#title.top = Cm(1.50)

	# set the liturgie:
	subtitle = slide.placeholders[1]
	for song_num in song_couplets.keys():
		title_text = "Lied "
		if song_num.isdigit() and int(song_num) <= 150: 
			title_text = "Psalm "
		title_text += song_num + ': '
		for c_idx in range(0, len(song_couplets[song_num])):
			title_text += song_couplets[song_num][c_idx] + (', ' if  c_idx < len(song_couplets[song_num])-1 else '')
		subtitle.text += title_text + '\n'
	for idx in range(0, len(scripture_fragments)):
		subtitle.text += '\nSchriftlezing {0}: {1}'.format(idx+1, scripture_fragments[idx])

# todo: read these parameters from the command-line

###################  command_line part ####################
def start_cmdline():
	voorganger = 'Ds. F. Schipper'
	datum_tekst = 'zondag 15 januari 2017'
	scripture_fragments = ['Marcus: 1-11',]
	titel_tekst = 'Welkom!'
	sub_titel_tekst = datum_tekst + '\nVoorganger: ' + voorganger
	create_ppt(voorganger, datum_tekst, scripture_fragments, titel_tekst, sub_titel_tekst)
	return

def create_ppt(voorganger, datum_tekst, scripture_fragments, titel_tekst, sub_titel_tekst):
	pptx_template_file = 'template.pptx'
	liedboek_file = 'liedboek.zip'
	
	if zipfile.is_zipfile(liedboek_file):
		zf = zipfile.ZipFile(liedboek_file, 'r')
	else:
		exit(liedboek_file + ' is not readable')
	
	filenamelist = sorted(zf.namelist()) # the zipfile contains files in unspecified order, so manual sort is necessary
	
	prs = create_pptx(pptx_template_file)
	create_title_slide(prs, titel_tekst, sub_titel_tekst)
	
	song_couplets = song_couplets2arr(filenamelist)
	create_index_slide(prs, song_couplets, scripture_fragments, datum_tekst)
	idx = 1
	for scripture_fragment in scripture_fragments:
		create_scripture_slide(prs, "Schriftlezing {0}: {1}".format(idx, scripture_fragment), "<tekst van {0} hier plakken>".format(scripture_fragment))
		idx += 1
	
	standaard_ochtenddienst_layout = ['titel', 'liturgie', 'lied1', 'Welkom en afkondigingen', 
				'Stil gebed\n-\nVotum en Groet', 'lied2', 'Lezing van Gods gebod',
				'lied3', 'Gebed om de opening van het Woord', 'Projectlied via de beamer',
				'Kinderen komen naar voren en gaan naar de kindernevendienst',
				'schriftlezing1', 'lied4', 'Verkondiging', 'lied5', 'Dankgebed',
				'Inzameling van de gaven', 'lied6', 'Zegen']
	standaard_avonddienst_layout = ['titel', 'liturgie', 'lied1', 'Stil gebed\n-\nVotum en Groet',
					'lied2', 'Gebed om de opening van het Woord', 'schriftlezing1',
					'lied4', 'Verkondiging', 'lied5', 'Geloofsbelijdenis', 'lied6',
					'Gebed', 'Inzameling van de gaven', 'lied6', 'Zegen']
	
	for dianaam in standaard_ochtenddienst_layout:
		create_intermediate_slide(prs, dianaam)
	
	
	for filename in filenamelist:
		if filename[-3:] == 'png':
			print 'processing img: {0}'.format(filename)
			song_img_data = zf.read(filename)
			img = Image.open(StringIO.StringIO(song_img_data))
			#img = Image.open(song_img_data)
			width, height = img.size
			img2 = img.crop((0, 150 , width, height))  #origineel komt binnen als 1600x1200 (haal van de bovenkant 150 px af
			# do not do any resizing here, but leave the original size and resizing using the height en width attributes of the shape (picture object), because this results in a sharper image
			img3 = StringIO.StringIO()
			img2.save(img3, format='PNG', quality=100)
	
			song_title = get_song_title_text(filename, song_couplets)
			create_song_slide(prs, song_title, img3)
	
	prs.save('result.pptx')
	print "powerpoint created..."
	zf.close()


def start_web():
	print "this is start_web"






###################  web part ####################
UPLOAD_FOLDER = '/tmp/pytest'
ALLOWED_EXTENSIONS = set(['zip'])

def allowed_file(filename):
	return '.' in filename and \
		filename.rsplit('.', 1)[1] in ALLOWED_EXTENSIONS

app = Flask(__name__)
app.secret_key = 'some_secret'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # max 16 MB


#@app.route('/')
#def hello_world():
#   return 'Hello, World!'

@app.route('/', methods=['GET', 'POST'])
def index():
	filenamelist = get_filenamelist()
	song_couplets = song_couplets2arr(filenamelist)
	liturgielijst = []
	for song, couplets in song_couplets.iteritems():
		#for couplet in couplets:  # todo: per couplet sorteren mogelijk maken...
		coupletstr = ', '.join(couplets)
		liturgielijst.append([song, coupletstr])  #'{0}: {1}'.format(song, coupletstr))
	#liturgielijst = song_couplets   #{ 'title': 'allard', 'Age': 7 }
	return render_template('index.html', name='test van Allard',liturgielijst=liturgielijst)


@app.route('/login', methods=['GET', 'POST'])
def login():
	error = None
	if request.method == 'POST':
		if request.form['username'] != 'admin' or \
			request.form['password'] != 'secret':
			error = 'Invalid credentials'
		else:
			flash('You were successfully logged in')
			return redirect(url_for('index'))
	return render_template('login.html', error=error)



@app.route('/upload', methods=['GET', 'POST'])
def upload_file():
	error=None
	if request.method == 'POST':
		# check if the post request has the file part
		if 'file' not in request.files:
			flash('No file part')
			return redirect(request.url)
		file = request.files['file']
		# if user does not select file, browser also
		# submit a empty part without filename
		if file.filename == '':
			flash('No selected file')
			return redirect(request.url)
		if file:
			if allowed_file(file.filename):
				filename = secure_filename(file.filename)
				file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
				#return redirect(url_for('upload_file', filename=filename))
				flash('upload=suc6')
				return redirect(url_for('upload_file', filename=filename))
			else:
				flash('Invalid filetype (only .zip is allowed)')

	return render_template('upload.html', introtekst='Upload nieuw bestand', errormsg=error)



@app.route('/pptx_summary', methods=['POST'])
def pptx_summary():
	if request.method == 'POST':
		finalliturgielijst = []
		# check if the post request has the file part
		if 'liedvolgorde' not in request.form:
			flash('Geen liederen gevonden')
		else:
			liedstr = request.form['liedvolgorde'].split(',')
			liturgietypestr = request.form['liturgietype']
			flash('Wel liederen gevonden {0}. Liturgietype={1}'.format(liedstr,liturgietypestr))

		# todo Allard: ga hier verder
		finalliturgielijst.append(['test3', 'test4'])
		return render_template('pptxsummary.html', finalliturgielijst=finalliturgielijst)


# todo:
@app.route('/pptx_save', methods=['POST', 'GET'])
def pptx_save():
	if request.method == 'POST':
		# check if the post request has the file part
		if 'liturgievolgorde' not in request.form:
			flash('Geen liederen gevonden')
		else:
			litstr = request.form['liturgievolgorde'].split(',')
			flash('Wel liederen gevonden {0}'.format(litstr))
	else:
		return render_template('saving.html', introtekst='Saving 5')

	return redirect(request.url)

############################ entry point (main) ####################
if __name__ == "__main__":
	arv = sys.argv[1:]
	if(len(arv)>0):
		if arv[0] == '-c':  # start command-line version
			start_cmdline()
		elif arv[0] == '-w':  # start the webserver version
			start_web()
		else:
			print "invalid argument"
	else:
		print "no arguments given! usage: \n  -c (start the command_line version)\n  -w (start the webserver version)"

