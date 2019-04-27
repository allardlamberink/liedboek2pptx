#!/usr/bin/env python

import sys
from threading import Thread
from time import sleep
import zipfile
import StringIO
from pptx import Presentation
from pptx.util import Cm
from PIL import Image
import os
from collections import OrderedDict


class CreatePPTXProcess(Thread):
    params_are_set = False
    total_file_count = 0
    files_processed_count = 0
    upload_path = None
    uploaded_zipfilename = None
    voorganger = None
    datum_tekst = None
    scripture_fragments = None
    titel_tekst = None
    sub_titel_tekst = None
    liedvolgorde = None
 

    def __init__(self, *args, **kwargs):
        Thread.__init__(self)
        self.files_processed_count = 0
        self.key = kwargs.get('file_uuid', None)


    def setparams(self, upload_path, uploaded_zipfilename, liedvolgorde, voorganger, datum_tekst, scripture_fragments, titel_tekst, sub_titel_tekst):
        self.upload_path = upload_path
        self.uploaded_zipfilename = uploaded_zipfilename
        self.voorganger = voorganger
        self.datum_tekst = datum_tekst
        self.scripture_fragments = scripture_fragments #['Mattheus 5: 1-15', 'John 3: 16']
        self.titel_tekst = titel_tekst
        self.sub_titel_tekst = sub_titel_tekst #+ '\nVoorganger: ' + voorganger
        self.liedvolgorde = liedvolgorde
        self.params_are_set = True


    # filename structuur binnen zip_obj:
    # meerdere coupletten:
    # projectie-111-muziek-couplet-1-1.png
    # maar 1 couplet:
    # projectie-425-muziek-2.png
    def song_couplets2arr(self, filenamelist):
        song_couplets = {}
        for filename in filenamelist:
            if filename[-3:] == 'png':
                title_text_arr = filename.split('-')
                if title_text_arr[3][0].isdigit() == False:  # meerdere coupletten
                    if title_text_arr[1] in song_couplets:
                        if title_text_arr[4] not in song_couplets[title_text_arr[1]]:
                            song_couplets[title_text_arr[1]].append(title_text_arr[4])
                    else:
                        song_couplets[title_text_arr[1]] = [title_text_arr[4]]
                else:  # 1 couplet
                    song_couplets[title_text_arr[1]] = ['1']
        return song_couplets


    # filename structuur binnen zip_obj:
    # meerdere coupletten:
    # projectie-111-muziek-couplet-1-1.png
    # maar 1 couplet:
    # projectie-425-muziek-2.png
    def sort_filenamelist(self, filenamelist, volgordelist):
        sorted_filenamelist = []
        for liednr in volgordelist:
            for filename in filenamelist:
                if filename.startswith('projectie-%s-' % liednr):
                    sorted_filenamelist.append(filename)
        return sorted_filenamelist


    def get_song_title_text(self, filename, song_couplets):
        title_text = 'Lied '
        title_text_arr = filename.split('-')
        if title_text_arr[1].isdigit() and int(title_text_arr[1]) <= 150:
            title_text = "Psalm "
        title_text += title_text_arr[1] + ': '
    
        #print len(song_couplets[title_text_arr[1]])
        for couplet in song_couplets[title_text_arr[1]]:
            if len(song_couplets[title_text_arr[1]])==1:
                title_text += ' [' + couplet + '] '
            elif couplet == title_text_arr[4]:
                title_text += ' [' + couplet + '] '
            else:
                title_text += ' ' + couplet + ' '
        return title_text


    def create_pptx(self, pptx_template_file):
        prs = Presentation(pptx_template_file)  # if filename is give, load the presentation
        prs.core_properties.author = "pptx-generator v0.1"
        prs.core_properties.title = "pptx-generator generator powerpoint for Hervgemb"
        return prs


    def create_title_slide(self, prs, titel_tekst, sub_titel_tekst):
        title_slide_layout = prs.slide_layouts[0]  # layout 0 = de startpagina
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = titel_tekst
        subtitle.text = sub_titel_tekst


    def create_song_slide(self, prs, song_title, song_img_data):
        song_slide_layout = prs.slide_layouts[1]  # layout 1 = titel + object
        slide = prs.slides.add_slide(song_slide_layout)
    
        # set title Psalm / Lied
        title = slide.shapes.title
        title.text = song_title
        #title.left = Cm(3.30)
        #title.top = Cm(1.50)
    
        # set song image (SongText + MusicNotes)
        left = Cm(3.29)
        top = Cm(1.94)
        #pic = slide.shapes.add_picture(song_img_data, left, top)
        img4= StringIO.StringIO(song_img_data.getvalue())
        pic = slide.shapes.add_picture(img4, left, top)
        #pic = slide.shapes.add_picture(song_img_data, left, top)
        pic.height=pic.height/3
        pic.width=pic.width/3


    def create_scripture_slide(self, prs, scripture_title, scripture_text):
        scripture_slide_layout = prs.slide_layouts[1]  # layout 1 = titel + object
        slide = prs.slides.add_slide(scripture_slide_layout)
    
        # set scripture title
        title = slide.shapes.title
        title.text = scripture_title
    
        # set scripture text
        subtitle = slide.placeholders[1]
        subtitle.top = Cm(1.94)
        subtitle.left = Cm(1.00)
        subtitle.width = prs.slide_width - subtitle.left - Cm(0.50)
        subtitle.height = prs.slide_height - subtitle.top - Cm(0.50)
        subtitle.text += scripture_text + '\n'


    def create_intermediate_slide(self, prs, tekst):
        interim_slide_layout = prs.slide_layouts[3]  # layout 3 = intermediate_slide
        slide = prs.slides.add_slide(interim_slide_layout)
    
        # set title
        title = slide.shapes.title
        title.text = tekst
        #title.left = Cm(3.30)
        #title.top = Cm(1.50)


    def create_index_slide(self, prs, song_couplets, scripture_fragments, datum_tekst):
        index_slide_layout = prs.slide_layouts[1]  # layout 1 = titel + object
        slide = prs.slides.add_slide(index_slide_layout)
    
        # set title Psalm / Lied
        title = slide.shapes.title
        title.text = 'Liturgie ' + datum_tekst
        #title.left = Cm(3.30)
        #title.top = Cm(1.50)
    
        # set the liturgie:
        subtitle = slide.placeholders[1]
        for song_num in song_couplets.keys():
            title_text = "Lied "
            if song_num.isdigit() and int(song_num) <= 150:
                title_text = "Psalm "
            title_text += song_num + ': '
            for c_idx in range(0, len(song_couplets[song_num])):
                title_text += song_couplets[song_num][c_idx] + (', ' if  c_idx < len(song_couplets[song_num])-1 else '')
            subtitle.text += title_text + '\n'
        for idx in range(0, len(scripture_fragments)):
            subtitle.text += '\nSchriftlezing {0}: {1}'.format(idx+1, scripture_fragments[idx])


    def get_zip_obj(self, zipfilename):
        if zipfile.is_zipfile(zipfilename):
            zip_obj = zipfile.ZipFile(zipfilename, 'r')
        else:
            exit(zipfilename + ' is not readable')
        return zip_obj


    def get_filenamelist(self, zip_obj):
        filenamelist = sorted(zip_obj.namelist()) # the zipfile contains files in unspecified order, so manual sort is necessary
        self.total_file_count = len(filenamelist)
        return filenamelist


    def create_ppt(self, uploaded_zipfilename, volgordelist, voorganger, datum_tekst, scripture_fragments, titel_tekst, sub_titel_tekst):
        pptx_template_file = 'template.pptx'
        zip_obj = self.get_zip_obj(uploaded_zipfilename)
        filenamelist = self.get_filenamelist(zip_obj)
        
        sorted_filenamelist = self.sort_filenamelist(filenamelist, volgordelist)
    
        prs = self.create_pptx(pptx_template_file)
        self.create_title_slide(prs, titel_tekst, sub_titel_tekst)
    
        song_couplets = self.song_couplets2arr(sorted_filenamelist)
        song_couplets_sorted = OrderedDict((k, song_couplets[k]) for k in volgordelist)


        self.create_index_slide(prs, song_couplets_sorted, scripture_fragments, datum_tekst)
        idx = 1
        for scripture_fragment in scripture_fragments:
            self.create_scripture_slide(prs, "Schriftlezing {0}: {1}".format(idx, scripture_fragment), "<tekst van {0} hier plakken>".format(scripture_fragment))
            idx += 1
    
        standaard_ochtenddienst_layout = ['titel', 'liturgie', 'lied1',
                    'Stil gebed\n-\nVotum en Groet', 'lied2', 'Lezing van Gods gebod',
                    'lied3', 'Gebed om de opening van het Woord', 'Kinderlied via de beamer',
                    'Kinderen komen naar voren en gaan naar de kindernevendienst',
                    'schriftlezing1', 'lied4', 'Verkondiging', 'lied5', 'Dankgebed',
                    'Inzameling van de gaven', 'lied6', 'Zegen']
        standaard_avonddienst_layout = ['titel', 'liturgie', 'lied1', 'Stil gebed\n-\nVotum en Groet',
                        'lied2', 'Gebed om de opening van het Woord', 'schriftlezing1',
                        'lied4', 'Verkondiging', 'lied5', 'Geloofsbelijdenis', 'lied6',
                        'Gebed', 'Inzameling van de gaven', 'lied6', 'Zegen']
    
        for dianaam in standaard_ochtenddienst_layout:
            self.create_intermediate_slide(prs, dianaam)
    
        #total_file_count = len(filenamelist)
        for filename in sorted_filenamelist:
            if filename[-3:] == 'png':
                print 'processing img: {0}'.format(filename)
                self.files_processed_count += 1
                song_img_data = zip_obj.read(filename)
                img = Image.open(StringIO.StringIO(song_img_data))
                #img = Image.open(song_img_data)
                width, height = img.size
                img2 = img.crop((0, 150 , width, height))  #origineel komt binnen als 1600x1200 (haal van de bovenkant 150 px af
                # do not do any resizing here, but leave the original size and resizing using the height en width attributes of the shape (picture object), because this results in a sharper image
                img3 = StringIO.StringIO()
                img2.save(img3, format='PNG', quality=100)
    
                song_title = self.get_song_title_text(filename, song_couplets_sorted)
                self.create_song_slide(prs, song_title, img3)

        file_with_path = os.path.join(self.upload_path, self.key)
        prs.save('%s.pptx' % file_with_path)
        print "powerpoint created... {0}".format(self.key)
        zip_obj.close()


    def run(self):
        if self.params_are_set and self.key:
            self.create_ppt(self.uploaded_zipfilename, self.liedvolgorde, self.voorganger, self.datum_tekst, self.scripture_fragments, self.titel_tekst, self.sub_titel_tekst)
        else:
            print "make sure all parameters and key are set"


    def percent_done(self):
        """Gets the current percent done for the thread."""
        if self.total_file_count != 0:
            return float(self.files_processed_count) / float(self.total_file_count) * 100.0
        else:
            return 1


    def get_progress(self):
        """Can be called at any time before, during or after thread
        execution, to get current progress."""
        return '%d files (%.2f%%)' % (self.files_processed_count, self.percent_done())


class CreatePPTXProcessShellRun(object):
    """Runs an instance of the thread with shell output / feedback."""
    
    def __init__(self, init_class=CreatePPTXProcess):
        self.init_class = init_class


    def __call__(self, *args, **kwargs):
        cxp = self.init_class(*args, **kwargs)

        print '%s threaded process beginning.' % cxp.__class__.__name__
        print '%d files will be processed. ' % cxp.total_file_count + 'Now beginning progress output.' 
        print cxp.get_progress()

        cxp.start()

        while cxp.is_alive() and cxp.files_processed_count < cxp.total_file_count:
            sleep(1.0)
            print cxp.get_progress()

        print '%s threaded process complete. Now exiting.' % cxp.__class__.__name__


if __name__ == '__main__':
    CreatePPTXProcessShellRun()()
